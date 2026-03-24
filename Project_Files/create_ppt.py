from pptx import Presentation
from pptx.util import Inches, Pt
import pandas as pd

prs = Presentation()

# Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Parkinson's Disease Detection"
subtitle.text = "Machine Learning Pipeline\nPerformance and Methodology\nDeveloped by Piyush Kadam (@piyushkadam96k)"

# Introduction Slide
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title_shape = slide.shapes.title
title_shape.text = "Introduction & Objectives"
tf = slide.placeholders[1].text_frame
tf.text = "Parkinson's Disease (PD) is a neurodegenerative disorder with significant motor and non-motor implications."
p = tf.add_paragraph()
p.text = "Early diagnosis is critical but challenging relying purely on clinical signs."
p = tf.add_paragraph()
p.text = "Voice degradation (tremors, changes in frequency/amplitude) presents as a powerful early biomarker."
p = tf.add_paragraph()
p.text = "Objective: To predict PD using machine learning applied to vocal features with >90% accuracy."

# The Dataset Slide
slide = prs.slides.add_slide(slide_layout)
title_shape = slide.shapes.title
title_shape.text = "The Dataset"
tf = slide.placeholders[1].text_frame
tf.text = "Source: UCI Machine Learning Repository"
p = tf.add_paragraph()
p.text = "Contains 195 biomedical voice measurements from 31 individuals (23 with PD)."
p = tf.add_paragraph()
p.text = "Key Features analyzed:"
p = tf.add_paragraph()
p.text = "- Fundamental Frequency: MDVP:Fo(Hz)"
p.level = 1
p = tf.add_paragraph()
p.text = "- Jitter & Shimmer: Variation in frequency/amplitude"
p.level = 1
p = tf.add_paragraph()
p.text = "- HNR, NHR: Noise to Harmonics ratios"
p.level = 1

# Methodology Slide
slide = prs.slides.add_slide(slide_layout)
title_shape = slide.shapes.title
title_shape.text = "Robust Methodology"
tf = slide.placeholders[1].text_frame
tf.text = "A strict machine learning pipeline was established to ensure real-world validity:"
p = tf.add_paragraph()
p.text = "1. Data Splitting: The dataset is split 80% Train, 20% Test before any transformations to prevent data leakage."
p = tf.add_paragraph()
p.text = "2. SMOTE Balancing: Synthetic Minor Oversampling Technique was applied ONLY on the training data to fix class imbalance."
p = tf.add_paragraph()
p.text = "3. Min-Max Scaling: Normalizes feature magnitudes without leaking test parameters."
p = tf.add_paragraph()
p.text = "4. GridSearchCV: Accelerated multi-core hyperparameter tuning to find optimal model configurations."

# Models Slide
slide = prs.slides.add_slide(slide_layout)
title_shape = slide.shapes.title
title_shape.text = "Models Deployed"
tf = slide.placeholders[1].text_frame
tf.text = "Several classifiers were trained and comprehensively cross-validated:"
p = tf.add_paragraph()
p.text = "1. Support Vector Machines (SVM) - Linear and RBF Kernels"
p = tf.add_paragraph()
p.text = "2. XGBoost - Gradient Boosting"
p = tf.add_paragraph()
p.text = "3. Random Forest - Ensemble Trees"
p = tf.add_paragraph()
p.text = "4. Decision Trees, Logistic Regression, KNN, and Naive Bayes"

# Results Slide (Table)
slide_layout = prs.slide_layouts[5] # Title only layout for table
slide = prs.slides.add_slide(slide_layout)
title_shape = slide.shapes.title
title_shape.text = "Pipeline Evaluation Metrics"

# Load data and add table
df = pd.read_csv('model_metrics_comparison.csv')
rows = min(8, len(df) + 1)
cols = 4 # Model, Accuracy, F1, Recall
left = Inches(0.5)
top = Inches(2.0)
width = Inches(9.0)
height = Inches(0.8 * rows)

table = slide.shapes.add_table(rows, cols, left, top, width, height).table
# Set column widths
table.columns[0].width = Inches(2.5)
table.columns[1].width = Inches(1.5)
table.columns[2].width = Inches(1.5)
table.columns[3].width = Inches(1.5)

# Set headers
headers = ['Classifier Model', 'Accuracy', 'F1-Score', 'Recall']
for i, header in enumerate(headers):
    table.cell(0, i).text = header

# Fill data
for i in range(len(df)):
    if i >= 7: break
    table.cell(i+1, 0).text = str(df.iloc[i]['Model'])
    table.cell(i+1, 1).text = f"{df.iloc[i]['Accuracy']*100:.1f}%"
    table.cell(i+1, 2).text = f"{df.iloc[i]['F1-Score']*100:.1f}%"
    table.cell(i+1, 3).text = f"{df.iloc[i]['Recall']*100:.1f}%"

# Conclusion Slide
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title_shape = slide.shapes.title
title_shape.text = "Conclusion"
tf = slide.placeholders[1].text_frame
tf.text = "Support Vector Machine (SVM) achieved the highest real-world diagnostic performance with ~92.3% Accuracy and 94.1% F1-Score."
p = tf.add_paragraph()
p.text = "XGBoost and Random Forest closely follow, demonstrating the strength of ensembles."
p = tf.add_paragraph()
p.text = "By correcting the data-leakage pipeline flaw, these metrics now represent a realistic benchmark for non-invasive, voice-based Parkinson's diagnostic tools."

prs.save("Parkinsons_Project_Presentation.pptx")
print("Presentation generated successfully!")
